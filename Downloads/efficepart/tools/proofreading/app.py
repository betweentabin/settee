from flask import Flask, request, jsonify, render_template, session
from werkzeug.utils import secure_filename
import os
import json
import uuid
from datetime import datetime
import spacy
import docx
import PyPDF2
import xlrd
from pptx import Presentation
import io
import tempfile
import google.generativeai as genai  # 正しいインポート方法
import re

# Flaskアプリケーションの初期化
app = Flask(__name__)
app.secret_key = os.environ.get('SECRET_KEY', 'development_key')
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 500 * 1024 * 1024  # 最大500MBのファイル
app.config['GOOGLE_API_KEY'] = 'AIzaSyAqWDwIwflmIjjXvO76RC5EzytII9hDT8U'  # Google API Key

# Gemini APIの初期化
genai.configure(api_key=app.config['GOOGLE_API_KEY'])
model = genai.GenerativeModel('gemini-1.5-flash')

# アップロードフォルダがなければ作成
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# SpaCyモデルの読み込み
try:
    nlp = spacy.load("ja_core_news_sm")
except OSError:
    print("SpaCyの日本語モデルをダウンロードしてください")
    print("実行コマンド: python -m spacy download ja_core_news_sm")
    exit(1)

# 校正履歴を保存する辞書
proofreading_history = {}

# 許可するファイル拡張子
ALLOWED_EXTENSIONS = {
    'txt', 'doc', 'docx', 'pdf', 'xls', 'xlsx', 
    'ppt', 'pptx', 'odt', 'ods', 'odp'
}

def allowed_file(filename):
    """アップロードされたファイルの拡張子が許可されているかチェックする関数"""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/')
def index():
    """トップページを表示"""
    return render_template('index.html')

# Google Gemini APIを使用して文章構成を行う関数
def generate_text_structure(prompt, max_tokens=1500, structure_type='detailed'):
    """
    Google Gemini APIを使用して文章構成を生成する
    
    Args:
        prompt: 文章構成のプロンプト
        max_tokens: 生成するトークンの最大数
        structure_type: 文章構成のタイプ（'detailed'または'simple'）
        
    Returns:
        生成された文章構成
    """
    # 文章構成のタイプに応じてプロンプトを変更
    if structure_type == 'simple':
        prompt_text = f"""以下のテーマについて簡潔な文章構成を作成してください。
        
テーマ: {prompt}

以下の形式で回答してください：
1. はじめに
2. 本論（3〜5つの主要ポイント）
3. 結論

各セクションの内容を簡潔に箇条書きで示してください。"""
    else:  # detailed
        prompt_text = f"""以下のテーマについて詳細な文章構成を作成してください。
        
テーマ: {prompt}

以下の形式で回答してください：
1. はじめに
   - 背景
   - 目的
   - 概要

2. 本論（複数の章に分けて）
   - 各章のタイトル
   - 各章で扱うべき重要なポイント（箇条書き）
   - 引用すべき情報や事例の提案

3. 結論
   - まとめ
   - 今後の展望

できるだけ具体的で、読者が理解しやすい構成にしてください。"""
    
    try:
        # コンテンツ生成
        response = model.generate_content(
            prompt_text,
            generation_config={
                "temperature": 0.7,
                "max_output_tokens": max_tokens,
            }
        )
        
        # 生成されたテキストを返す
        return response.text
    except Exception as e:
        print(f"API呼び出しエラー: {e}")
        return f"エラーが発生しました: {str(e)}"

@app.route('/generate-structure', methods=['POST'])
def generate_structure():
    """文章構成を生成するエンドポイント"""
    data = request.get_json()
    
    if not data or 'prompt' not in data:
        return jsonify({'error': 'プロンプトが必要です'}), 400
    
    prompt = data['prompt']
    max_tokens = data.get('max_tokens', 1500)
    
    # 文章構成のタイプ（詳細または簡潔）
    structure_type = data.get('type', 'detailed')
    
    # 文章構成を生成
    structure = generate_text_structure(prompt, max_tokens, structure_type)
    
    # 履歴に追加
    structure_id = str(uuid.uuid4())
    timestamp = datetime.now().isoformat()
    
    if 'structures' not in proofreading_history:
        proofreading_history['structures'] = {}
    
    proofreading_history['structures'][structure_id] = {
        'id': structure_id,
        'prompt': prompt,
        'structure': structure,
        'type': structure_type,
        'timestamp': timestamp
    }
    
    return jsonify({
        'id': structure_id,
        'prompt': prompt,
        'structure': structure,
        'type': structure_type,
        'timestamp': timestamp
    })

@app.route('/proofread', methods=['POST'])
def proofread():
    """テキストを校正するAPIエンドポイント"""
    data = request.get_json()
    
    if not data or 'text' not in data:
        return jsonify({'error': '校正するテキストが必要です'}), 400
    
    text = data['text']
    
    try:
        # 校正処理
        correction_result = perform_proofreading(text)
        
        # 結果オブジェクトの作成
        result = {
            'original': text,
            'corrected': text,  # 実際には修正されたテキスト
            'suggestions': correction_result
        }
        
        # 文章構成の提案を取得（オプション）
        if len(text) > 100 and 'suggest_structure' in data and data['suggest_structure']:
            try:
                structure = generate_text_structure(text)
                result['structure_suggestion'] = structure
            except Exception as e:
                print(f"文章構成生成エラー: {e}")
                # 構成生成に失敗しても校正結果は返す
        
        # 履歴に追加
        item_id = str(uuid.uuid4())
        timestamp = datetime.now().isoformat()
        
        # 履歴構造の初期化
        if 'items' not in proofreading_history:
            proofreading_history['items'] = {}
        
        # 履歴に保存
        proofreading_history['items'][item_id] = {
            'id': item_id,
            'original': text,
            'corrected': text,
            'suggestions': correction_result,
            'timestamp': timestamp
        }
        
        return jsonify(result)
        
    except Exception as e:
        print(f"校正処理中にエラー: {str(e)}")
        error_type = type(e).__name__
        
        # 適切なエラーメッセージを準備
        if "API" in str(e):
            error_message = "APIエラー: サーバー側の問題が発生しました。しばらく経ってからもう一度お試しください。"
        else:
            error_message = f"校正処理中にエラーが発生しました: {str(e)}"
        
        return jsonify({
            'error': error_message,
            'original': text,
            'suggestions': []
        }), 500

@app.route('/history')
def get_history():
    """校正履歴を取得するエンドポイント"""
    # 履歴を日付の新しい順にソート
    sorted_history = []
    if 'items' in proofreading_history:
        sorted_history = sorted(
            proofreading_history['items'].values(),
            key=lambda x: x.get('timestamp', ''),
            reverse=True
        )
    
    # 文章構成履歴
    structures = {}
    if 'structures' in proofreading_history:
        structures = proofreading_history['structures']
    
    return jsonify({
        'history': sorted_history,
        'structures': structures
    })

def extract_text_from_file(file_path, file_extension):
    """各種ファイル形式からテキストを抽出する関数"""
    try:
        print(f"テキスト抽出開始: {file_path}, 拡張子: {file_extension}")
        
        if file_extension == 'txt':
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
                    print(f"テキストファイルから {len(text)} 文字を抽出しました")
                    return text
            except UnicodeDecodeError:
                # UTF-8で失敗した場合、Shift-JISで試行
                with open(file_path, 'r', encoding='shift_jis') as f:
                    text = f.read()
                    print(f"テキストファイルから {len(text)} 文字を抽出しました（Shift-JIS）")
                    return text
        
        elif file_extension in ['doc', 'docx']:
            try:
                doc = docx.Document(file_path)
                text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
                print(f"Wordファイルから {len(text)} 文字を抽出しました")
                return text
            except Exception as e:
                print(f"Wordファイルの処理中にエラー: {str(e)}")
                raise
        
        elif file_extension == 'pdf':
            try:
                text = []
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        text.append(page.extract_text())
                result = '\n'.join(text)
                print(f"PDFファイルから {len(result)} 文字を抽出しました")
                return result
            except Exception as e:
                print(f"PDFファイルの処理中にエラー: {str(e)}")
                raise
        
        elif file_extension in ['xls', 'xlsx']:
            try:
                workbook = xlrd.open_workbook(file_path)
                text = []
                for sheet in workbook.sheets():
                    for row in range(sheet.nrows):
                        text.append(' '.join(str(cell.value) for cell in sheet.row(row)))
                result = '\n'.join(text)
                print(f"Excelファイルから {len(result)} 文字を抽出しました")
                return result
            except Exception as e:
                print(f"Excelファイルの処理中にエラー: {str(e)}")
                raise
        
        elif file_extension in ['ppt', 'pptx']:
            try:
                prs = Presentation(file_path)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                result = '\n'.join(text)
                print(f"PowerPointファイルから {len(result)} 文字を抽出しました")
                return result
            except Exception as e:
                print(f"PowerPointファイルの処理中にエラー: {str(e)}")
                raise
        
        else:
            raise ValueError(f'未対応のファイル形式です: {file_extension}')
    
    except Exception as e:
        print(f"テキスト抽出中にエラーが発生しました: {str(e)}")
        raise Exception(f'テキスト抽出中にエラーが発生しました: {str(e)}')

@app.route('/upload', methods=['POST'])
def upload_file():
    """ファイルをアップロードして校正するエンドポイント"""
    print("ファイルアップロードリクエストを受信しました")
    
    if 'file' not in request.files:
        print("ファイルがリクエストに含まれていません")
        return jsonify({'error': 'ファイルがアップロードされていません'}), 400
    
    file = request.files['file']
    print(f"アップロードされたファイル: {file.filename}")
    
    if file.filename == '':
        print("ファイル名が空です")
        return jsonify({'error': 'ファイルが選択されていません'}), 400
    
    # ファイル名から拡張子を取得（日本語ファイル名に対応）
    original_filename = file.filename
    if '.' not in original_filename:
        print("ファイルに拡張子がありません")
        return jsonify({'error': 'ファイルに拡張子がありません。拡張子付きのファイルを選択してください。'}), 400
    
    # 拡張子を取得
    file_extension = original_filename.rsplit('.', 1)[1].lower()
    print(f"ファイル拡張子: {file_extension}")
    
    if file_extension not in ALLOWED_EXTENSIONS:
        print(f"許可されていないファイル形式です: {file_extension}")
        return jsonify({'error': f'許可されていないファイル形式です: {file_extension}。対応形式: {", ".join(ALLOWED_EXTENSIONS)}'}), 400
    
    # 安全なファイル名を生成（日本語を保持）
    filename = secure_filename(original_filename)
    if not filename:  # secure_filenameが空文字を返した場合
        filename = f"upload_{uuid.uuid4().hex[:8]}.{file_extension}"
    
    temp_dir = None
    temp_path = None
    
    try:
        # 一時ファイルとして保存
        temp_dir = tempfile.mkdtemp()
        temp_path = os.path.join(temp_dir, filename)
        
        # ファイルを保存
        file.save(temp_path)
        print(f"一時ファイルを保存しました: {temp_path}")
        
        # ファイルからテキストを抽出
        print("テキスト抽出を開始します")
        extracted_text = extract_text_from_file(temp_path, file_extension)
        print(f"テキスト抽出完了: {len(extracted_text)} 文字")
        
        if not extracted_text:
            print("テキストの抽出に失敗しました")
            return jsonify({'error': 'テキストの抽出に失敗しました。ファイルの内容を確認してください。'}), 500
        
        # 校正処理
        print("校正処理を開始します")
        corrections = perform_proofreading(extracted_text)
        print(f"校正処理完了: {len(corrections)} 件の校正提案")
        
        # 校正結果を生成
        result = {
            'original': extracted_text,
            'corrected': extracted_text,  # 実際には修正されたテキスト
            'suggestions': corrections,
            'filename': original_filename  # 元のファイル名を使用
        }
        
        # 履歴に追加
        item_id = str(uuid.uuid4())
        timestamp = datetime.now().isoformat()
        
        # 履歴構造の初期化
        if 'items' not in proofreading_history:
            proofreading_history['items'] = {}
        
        # 履歴に保存
        proofreading_history['items'][item_id] = {
            'id': item_id,
            'original': extracted_text,
            'corrected': extracted_text,  # 実際には修正されたテキスト
            'suggestions': corrections,
            'filename': original_filename,  # 元のファイル名を使用
            'timestamp': timestamp
        }
        
        print("校正結果を返します")
        return jsonify(result)
        
    except Exception as e:
        print(f"ファイル処理エラー: {str(e)}")
        return jsonify({'error': f'ファイル処理エラー: {str(e)}'}), 500
        
    finally:
        # 一時ファイルとディレクトリを削除
        if temp_path and os.path.exists(temp_path):
            try:
                os.remove(temp_path)
                print("一時ファイルを削除しました")
            except Exception as e:
                print(f"一時ファイルの削除中にエラーが発生しました: {str(e)}")
        
        if temp_dir and os.path.exists(temp_dir):
            try:
                os.rmdir(temp_dir)
                print("一時ディレクトリを削除しました")
            except Exception as e:
                print(f"一時ディレクトリの削除中にエラーが発生しました: {str(e)}")

def perform_proofreading(text):
    """文章校正のコア処理を行う関数"""
    # SpaCyで文章を解析
    doc = nlp(text)
    
    suggestions = []
    
    # 1. SpaCyを使った基本的なチェック
    # 記号の前後のスペースチェック
    pos = 0
    for token in doc:
        if token.is_punct and token.text in ["、", "。"]:
            # 記号の前にスペースがある場合
            if pos > 0 and text[pos-1:pos] == " ":
                suggestions.append({
                    "position": pos-1,
                    "length": 2,
                    "suggestion": token.text,
                    "reason": "句読点の前にスペースは不要です",
                    "color": "blue"
                })
        pos += len(token.text)
    
    # 敬語表現のチェック
    for i, token in enumerate(doc):
        if token.lemma_ == "です" and i > 0 and doc[i-1].pos_ == "VERB":
            suggestions.append({
                "position": token.idx - len(doc[i-1].text),
                "length": len(doc[i-1].text) + len(token.text),
                "suggestion": f"{doc[i-1].text}ます",
                "reason": "「〜です」より「〜ます」の方が適切です",
                "color": "yellow"
            })
    
    # 重複表現のチェック
    for i in range(len(doc) - 1):
        if doc[i].text == doc[i+1].text and not doc[i].is_punct:
            suggestions.append({
                "position": doc[i].idx,
                "length": len(doc[i].text) * 2,
                "suggestion": doc[i].text,
                "reason": "単語が重複しています",
                "color": "red"
            })
    
    # 2. 追加の文章構成チェック
    # 長すぎる文のチェック
    sentences = [sent.text for sent in doc.sents]
    for sent in sentences:
        if len(sent) > 100:  # 100文字以上の文を長すぎると判断
            sent_start = text.find(sent)
            if sent_start >= 0:
                suggestions.append({
                    "position": sent_start,
                    "length": len(sent),
                    "suggestion": sent,
                    "reason": "文が長すぎます。複数の短い文に分けることを検討してください",
                    "color": "yellow"
                })
    
    # 接続詞の過剰使用チェック
    conjunctions = {}
    for token in doc:
        if token.pos_ == "CCONJ" or token.pos_ == "SCONJ":  # 接続詞
            if token.text in conjunctions:
                conjunctions[token.text] += 1
            else:
                conjunctions[token.text] = 1
    
    for conj, count in conjunctions.items():
        if count > 3:  # 同じ接続詞が3回以上使用されている場合
            positions = []
            pos = 0
            for token in doc:
                if token.text == conj:
                    positions.append(pos)
                pos += len(token.text)
            
            # 最初の位置は警告しない（最初の使用は問題ない）
            for pos in positions[2:]:  # 3回目以降の使用に警告
                suggestions.append({
                    "position": pos,
                    "length": len(conj),
                    "suggestion": f"別の表現（例：「{get_alternative_conjunction(conj)}」）",
                    "reason": f"「{conj}」が多用されています",
                    "color": "yellow"
                })
    
    # 文体の一貫性チェック
    # デス・マス調とダ・デアル調の混在をチェック
    has_desu_masu = False
    has_da_dearu = False
    
    for token in doc:
        if token.lemma_ in ["です", "ます"]:
            has_desu_masu = True
        elif token.lemma_ in ["だ", "である"]:
            has_da_dearu = True
    
    if has_desu_masu and has_da_dearu:
        # 文体が混在している場合
        suggestions.append({
            "position": 0,
            "length": min(50, len(text)),  # 最初の50文字を範囲とする
            "suggestion": text[:min(50, len(text))],
            "reason": "デス・マス調とダ・デアル調が混在しています。文体を統一してください",
            "color": "blue"
        })
    
    # 3. 日本語特有の文法チェック
    # 「は」と「が」の使い方
    check_wa_ga_usage(doc, suggestions)
    
    # 助詞の重複チェック
    check_duplicate_particles(doc, suggestions)
    
    # 主語と述語の対応チェック
    check_subject_predicate_correspondence(doc, suggestions)
    
    # 漢字とひらがなの適切な使用
    check_kanji_hiragana_usage(doc, suggestions, text)
    
    # 読点の過剰・過少チェック
    check_comma_usage(doc, suggestions, text)
    
    # 4. Google APIを使用した高度な校正
    try:
        print("Google APIを使用した校正を開始")
        api_suggestions = get_api_suggestions(text)
        if api_suggestions:
            print(f"API校正: {len(api_suggestions)}件の提案を受信")
            # APIの提案が既存の提案と重複している場合は追加しない
            for api_sugg in api_suggestions:
                # 重複チェック
                is_duplicate = False
                for existing_sugg in suggestions:
                    if abs(existing_sugg["position"] - api_sugg["position"]) < 5 and existing_sugg["suggestion"] == api_sugg["suggestion"]:
                        is_duplicate = True
                        break
                
                if not is_duplicate:
                    suggestions.append(api_sugg)
        else:
            print("API校正: 有効な提案がありませんでした")
    except Exception as e:
        print(f"Google API呼び出しエラー: {e}")
        # エラーをコンソールに出力するだけで、処理は続行
    
    print(f"校正完了: 合計{len(suggestions)}件の提案")
    return suggestions

def check_wa_ga_usage(doc, suggestions):
    """「は」と「が」の使い方をチェックする"""
    # 文脈に応じた「は」と「が」の使い分けチェック
    for i, token in enumerate(doc):
        if token.text == "は" and i > 0 and i < len(doc) - 1:
            # 「〜には〜」という形で、強調や対比の意味がない場合
            if doc[i-1].text.endswith("に"):
                suggestions.append({
                    "position": token.idx,
                    "length": 1,
                    "suggestion": "が",
                    "reason": "この文脈では「は」より「が」の方が自然です",
                    "color": "yellow"
                })
        
        elif token.text == "が" and i > 0 and i < len(doc) - 1:
            # 主題を示す場合は「は」を使う
            if doc[i-1].pos_ == "NOUN" and i > 1 and doc[i-2].pos_ != "NOUN":
                if "主題" in [t.text for t in doc[i-3:i]]:
                    suggestions.append({
                        "position": token.idx,
                        "length": 1,
                        "suggestion": "は",
                        "reason": "主題を示す場合は「が」より「は」の方が適切です",
                        "color": "yellow"
                    })

def check_duplicate_particles(doc, suggestions):
    """助詞の重複をチェックする"""
    for i in range(len(doc) - 2):
        if doc[i].pos_ == "ADP" and doc[i+1].pos_ != "PUNCT" and doc[i+2].pos_ == "ADP":
            # 助詞が短い間隔で連続している場合
            suggestions.append({
                "position": doc[i].idx,
                "length": doc[i+2].idx + len(doc[i+2].text) - doc[i].idx,
                "suggestion": f"{doc[i].text}～{doc[i+2].text}",
                "reason": "助詞が短い間隔で連続しています。文を簡略化することを検討してください",
                "color": "yellow"
            })

def check_subject_predicate_correspondence(doc, suggestions):
    """主語と述語の対応をチェックする"""
    # 簡易的なチェック（完全な解析には構文解析が必要）
    sentences = list(doc.sents)
    for sent in sentences:
        has_subject = False
        has_predicate = False
        
        for token in sent:
            # 主語の存在チェック（「〜は」「〜が」）
            if token.text in ["は", "が"] and token.i > 0 and token.i < len(doc) - 1:
                if doc[token.i-1].pos_ == "NOUN":
                    has_subject = True
            
            # 述語の存在チェック（動詞、形容詞など）
            if token.pos_ in ["VERB", "ADJ"] and token.i == len(sent) - 1:
                has_predicate = True
        
        # 長い文で主語か述語がない場合
        if len(sent.text) > 30 and (not has_subject or not has_predicate):
            suggestions.append({
                "position": sent.start_char,
                "length": len(sent.text),
                "suggestion": sent.text,
                "reason": "主語と述語の対応が不明確です。文の構造を見直してください",
                "color": "blue"
            })

def check_kanji_hiragana_usage(doc, suggestions, text):
    """漢字とひらがなの適切な使用をチェックする"""
    # 一般的に漢字で書くべき語彙のチェック
    kanji_words = {
        "わたし": "私",
        "おおきい": "大きい",
        "ちいさい": "小さい",
        "すこし": "少し",
        "たくさん": "沢山",
        "とき": "時",
        "ひと": "人",
        "もの": "物",
        "ところ": "所",
        "こと": "事"
    }
    
    for word, kanji in kanji_words.items():
        index = text.find(word)
        while index >= 0:
            # 単語の前後が他の文字の一部でないことを確認
            is_independent = True
            if index > 0 and text[index-1].isalpha():
                is_independent = False
            if index + len(word) < len(text) and text[index + len(word)].isalpha():
                is_independent = False
            
            if is_independent:
                suggestions.append({
                    "position": index,
                    "length": len(word),
                    "suggestion": kanji,
                    "reason": f"一般的には「{word}」より漢字の「{kanji}」を使います",
                    "color": "yellow"
                })
            
            index = text.find(word, index + 1)

def check_comma_usage(doc, suggestions, text):
    """読点の使用をチェックする"""
    # 長い文で読点が少ない場合
    sentences = list(doc.sents)
    for sent in sentences:
        if len(sent.text) > 50:  # 50文字以上の長い文
            commas = [t for t in sent if t.text == "、"]
            if len(commas) == 0:  # 読点がない
                # 挿入候補位置を探す（接続詞の後など）
                for token in sent:
                    if token.pos_ in ["CCONJ", "SCONJ"] and token.i < len(sent) - 1:
                        suggestions.append({
                            "position": token.idx + len(token.text),
                            "length": 0,
                            "suggestion": "、",
                            "reason": "長い文では読点を使うと読みやすくなります",
                            "color": "yellow"
                        })
                        break

def get_alternative_conjunction(conj):
    """接続詞の代替案を提供する"""
    alternatives = {
        "しかし": "ですが、けれども、一方",
        "そして": "また、それから、加えて",
        "だから": "それゆえ、したがって、そのため",
        "ただし": "ただ、もっとも、とはいえ",
        "また": "さらに、加えて、それから",
        "それから": "その後、続いて、そして",
        "なお": "ちなみに、ところで、さらに",
        "ところで": "さて、そういえば、話は変わるが",
        "たとえば": "例えば、一例として、具体的には",
        "すなわち": "つまり、要するに、言い換えれば"
    }
    
    return alternatives.get(conj, "別の接続詞")

def get_api_suggestions(text):
    """Google Gemini APIを使用して文章の校正提案を取得する"""
    try:
        # 長すぎるテキストの場合は分割（API制限対策）
        if len(text) > 30000:
            print(f"テキストが長すぎるため、最初の30000文字だけを処理します: {len(text)}文字")
            text = text[:30000]

        # セーフティー設定
        safety_settings = [
            {
                "category": "HARM_CATEGORY_HARASSMENT",
                "threshold": "BLOCK_NONE",
            },
            {
                "category": "HARM_CATEGORY_HATE_SPEECH",
                "threshold": "BLOCK_NONE",
            },
            {
                "category": "HARM_CATEGORY_SEXUALLY_EXPLICIT",
                "threshold": "BLOCK_NONE",
            },
            {
                "category": "HARM_CATEGORY_DANGEROUS_CONTENT",
                "threshold": "BLOCK_NONE",
            },
        ]

        # プロンプトの設定
        prompt = f"""あなたは日本語文章校正の専門家です。
以下の文章を日本語の観点から校正し、次の点に注意してください：

1. 誤字脱字
2. 文法的な誤り（助詞の使い方など）
3. 不自然な表現や冗長な表現
4. 敬語・丁寧語の適切な使用
5. 文体の一貫性
6. 読点と句点の適切な配置
7. 同じ言葉の不必要な繰り返し

入力文章:
{text}

レスポンスは必ず以下のJSON形式で返してください:
[
  {{
    "position": 0,
    "length": 1,
    "suggestion": "修正後の文字列",
    "reason": "修正理由",
    "color": "red"
  }}
]

注意事項：
・必ずJSON配列を返してください
・配列が空の場合は [] を返してください
・説明文やマークダウン記法は使用しないでください
・positionは0から始まる整数値です
・lengthは1以上の整数値です
・colorは "red"（誤字脱字）、"blue"（文法ミス）、"yellow"（表現改善）のいずれかです"""

        print(f"APIリクエスト送信: プロンプト長さ {len(prompt)} 文字")
        
        # コンテンツ生成
        try:
            # APIタイムアウト設定の追加
            response = model.generate_content(
                prompt,
                generation_config={
                    "temperature": 0.1,
                    "max_output_tokens": 2048,
                    "top_p": 0.95,
                    "top_k": 40
                },
                safety_settings=safety_settings
            )
            
            # 生成されたテキストを取得
            generated_text = response.text.strip()
            
            # デバッグ用にレスポンスを出力
            print(f"APIレスポンス: 長さ {len(generated_text)} 文字")
            if len(generated_text) > 0:
                print(f"レスポンス先頭: {generated_text[:100]}")
            
            try:
                # 直接JSONとしてパースを試みる
                api_suggestions = json.loads(generated_text)
                
                # リスト形式かチェック
                if isinstance(api_suggestions, list):
                    # 各提案の形式を検証
                    valid_suggestions = []
                    for suggestion in api_suggestions:
                        if validate_suggestion(suggestion):
                            # 数値型の確認と変換
                            suggestion["position"] = int(suggestion["position"])
                            suggestion["length"] = int(suggestion["length"])
                            valid_suggestions.append(suggestion)
                    
                    print(f"有効な提案数: {len(valid_suggestions)}")
                    return valid_suggestions
                else:
                    print(f"APIレスポンスがリスト形式ではありません: {type(api_suggestions)}")
                    return []
                
            except json.JSONDecodeError as e:
                print(f"JSONパースエラー: {e}")
                # テキストのクリーニングを試みる
                print("レスポンスのクリーニングを試みます")
                
                # マークダウンのコードブロックを削除
                cleaned_text = re.sub(r'```(?:json)?\s*|\s*```\s*', '', generated_text)
                # 前後の空白を削除
                cleaned_text = cleaned_text.strip()
                
                print(f"クリーニング後: {cleaned_text[:100]}")
                
                try:
                    api_suggestions = json.loads(cleaned_text)
                    if isinstance(api_suggestions, list):
                        valid_suggestions = []
                        for suggestion in api_suggestions:
                            if validate_suggestion(suggestion):
                                # 数値型の確認と変換
                                suggestion["position"] = int(suggestion["position"])
                                suggestion["length"] = int(suggestion["length"])
                                valid_suggestions.append(suggestion)
                        
                        print(f"クリーニング後の有効な提案数: {len(valid_suggestions)}")
                        return valid_suggestions
                    else:
                        print("クリーニング後もリスト形式ではありません")
                        return []
                except json.JSONDecodeError as e2:
                    print(f"クリーニング後もJSONパースに失敗しました: {e2}")
                    return []
                
        except Exception as api_error:
            error_type = type(api_error).__name__
            error_msg = str(api_error)
            print(f"API呼び出しエラー [{error_type}]: {error_msg}")
            
            # 詳細なエラー情報を取得
            if hasattr(api_error, 'status_code'):
                print(f"ステータスコード: {api_error.status_code}")
            if hasattr(api_error, 'response'):
                print(f"レスポンス情報: {repr(api_error.response)}")
            if hasattr(api_error, 'details'):
                print(f"詳細: {api_error.details}")
                
            # APIエラーの種類に応じた処理
            if "quota" in error_msg.lower():
                print("API使用量の制限に達しました")
            elif "permission" in error_msg.lower():
                print("API許可エラー - APIキーを確認してください")
            elif "time" in error_msg.lower():
                print("APIタイムアウトエラー")
            
            return []
            
    except Exception as e:
        error_type = type(e).__name__
        print(f"予期せぬエラー [{error_type}]: {str(e)}")
        import traceback
        traceback.print_exc()
        return []

def validate_suggestion(suggestion):
    """提案の形式が正しいかを検証する"""
    required_keys = {"position", "length", "suggestion", "reason", "color"}
    if not all(key in suggestion for key in required_keys):
        return False
    
    if not isinstance(suggestion["position"], (int, str)) or (isinstance(suggestion["position"], int) and suggestion["position"] < 0):
        return False
    
    if not isinstance(suggestion["length"], (int, str)) or (isinstance(suggestion["length"], int) and suggestion["length"] < 1):
        return False
    
    if not isinstance(suggestion["suggestion"], str):
        return False
    
    if not isinstance(suggestion["reason"], str):
        return False
    
    if suggestion["color"] not in ["red", "blue", "yellow"]:
        return False
    
    return True

def create_error_suggestion(text, error_message):
    """エラー時のデフォルト提案を生成する"""
    return [{
        "position": 0,
        "length": len(text),
        "suggestion": text,
        "reason": error_message,
        "color": "yellow"
    }]

# APIレスポンス処理部分の修正
def parse_api_response(response_text):
    try:
        # 空の応答をチェック
        if not response_text or response_text.isspace():
            print("空のAPIレスポンスを受信しました")
            return []
            
        # レスポンスのデバッグ出力
        print(f"APIレスポンス: {response_text[:100]}...")
        
        # JSONとして解析
        import json
        data = json.loads(response_text)
        return data
    except json.JSONDecodeError as e:
        print(f"JSONパースエラー: {e}")
        print(f"問題のレスポンス: {response_text[:200]}")
        
        # 空のリストを返す（エラー回避）
        return []

if __name__ == '__main__':
    app.run(debug=True) 